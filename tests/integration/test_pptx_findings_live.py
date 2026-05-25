"""Live integration test — PPTX + K7 FindingsAgent composition.

Until this test landed, PPTX was zero-coverage in the agent suite:
every K-series live test sourced its document from a real DOCX NDA,
and the runs/INDEX.jsonl audit trail had no record of the agent
ever seeing a PowerPoint file. That's a real gap for a regulated
audience — board decks, M&A pitch books, and earnings prep slides
are some of the densest live targets for an "exhaustive findings"
agent, and skipping them in the audit trail means we can't claim
the surface works end-to-end on the document types our users hand
us most often.

Pipeline exercised:
1. Generate a deterministic 15-slide board deck in-memory via
   ``python-pptx`` (no external fixture dependency).
2. Parse it via ``kaos_office.parse_pptx`` → ``ContentDocument``.
3. Store it as a VFS artifact via
   ``kaos_content.artifacts.store_document``.
4. Run :class:`~kaos_agents.tools.findings.AgentFindingsTool` with
   ``select_by="token"`` + ``selector_arg="cyber"``.
5. Assert the synthesized answer recovers the known mitigation
   phrase ("multi-factor authentication" / "penetration testing")
   that we plant on a specific slide, cites at least one
   finding_id, and that cost stays inside a sane envelope.

Skips cleanly when ``python-pptx`` or ``ANTHROPIC_API_KEY`` is
absent — the auto-recorder fixture in conftest.py captures the
LLM call traces under ``tests/integration/runs/<date>/``.
"""

from __future__ import annotations

import os
import tempfile
from pathlib import Path
from typing import Any

import pytest

from tests.integration._models import critic_model

pptx = pytest.importorskip(
    "pptx",
    reason=(
        "python-pptx not installed; install kaos-office[pptx] or `uv add --group dev python-pptx`"
    ),
)


requires_anthropic = pytest.mark.skipif(
    "ANTHROPIC_API_KEY" not in os.environ,
    reason="ANTHROPIC_API_KEY missing",
)


# ---------------------------------------------------------------------------
# Fixture deck — deterministic, in-memory
# ---------------------------------------------------------------------------


# The known-good mitigation we plant on slide 8. Asserted on in the
# synthesis step: at least one of these substrings must appear in
# the final answer, proving the pipeline actually surfaced the
# planted sentence rather than hallucinating something plausible.
PLANTED_MITIGATION = "multi-factor authentication and quarterly penetration testing"
PLANTED_SLIDE_NUMBER = 8

# Deterministic content for a fake board meeting deck. Slides are
# topical (financials, M&A, legal, HR, cyber, regulatory, ...) so
# the token selector "cyber" hits exactly one slide; this gives us
# a small, predictable enumerated set that's still big enough to
# exercise both filter pass-and-drop branches.
SLIDES: list[tuple[str, list[str]]] = [
    (
        "Q3 2025 Financial Review",
        [
            "Revenue grew 14% YoY to $312M.",
            "Operating margin expanded to 23.4% on cost discipline.",
            "Free cash flow conversion held at 89%.",
        ],
    ),
    (
        "Segment Performance",
        [
            "North America posted 11% growth driven by mid-market.",
            "EMEA softened with 3% growth, FX-adjusted.",
            "APAC accelerated to 28% on enterprise wins.",
        ],
    ),
    (
        "M&A Pipeline",
        [
            "Three active diligence streams: Project Aurora, Beacon, Cypress.",
            "Aurora term sheet expected late Q4 2025.",
            "Beacon paused pending antitrust review timing.",
        ],
    ),
    (
        "Legal Update — Active Matters",
        [
            "Two active commercial disputes, both in pre-trial discovery.",
            "Patent matter Acme v. Subsidiary scheduled for mediation in February.",
            "No new regulatory inquiries received this quarter.",
        ],
    ),
    (
        "Regulatory Calendar",
        [
            "10-K filing planned for February 18, 2026.",
            "SOX 404 control testing completed without exceptions.",
            "New EU AI Act obligations assessed; gap analysis underway.",
        ],
    ),
    (
        "Headcount and Hiring",
        [
            "Net headcount grew 4% to 2,140 employees.",
            "Engineering hiring weighted 60% to platform, 40% to AI.",
            "Voluntary attrition held at 9.1%, below industry benchmark.",
        ],
    ),
    (
        "Compensation Plan Calibration",
        [
            "Executive LTI mix shifts to 60% PSU, 40% RSU for FY2026.",
            "Salary structure refresh effective April 1, 2026.",
            "Stock plan share pool sufficient through FY2027.",
        ],
    ),
    # Slide 8 — the planted cyber risk slide. The selector
    # 'cyber' will pin every sentence on this slide, and the
    # mitigation phrase must surface in the synthesized answer.
    (
        "Cyber Risk Mitigation",
        [
            "Top cyber risk this quarter is credential stuffing against partner SSO.",
            (f"Board-approved cyber mitigation is {PLANTED_MITIGATION} across all admin systems."),
            "Tabletop cyber incident exercise completed in October; gaps tracked to closure by Q1.",
        ],
    ),
    (
        "Product Roadmap Highlights",
        [
            "AI assistant GA target shifted from January to March 2026.",
            "Two new connectors shipping in Q4: ServiceNow and Workday.",
            "Mobile parity initiative funded through Q2 2026.",
        ],
    ),
    (
        "Customer Success Metrics",
        [
            "Net retention rate held at 117% on enterprise cohort.",
            "Logo retention at 94%, down 1 point QoQ from a single large churn.",
            "NPS rebounded to 52 from 48 last quarter.",
        ],
    ),
    (
        "Capital Allocation",
        [
            "$120M share repurchase authorized through end of FY2026.",
            "No dividend changes planned this fiscal year.",
            "Debt refinancing window targeted for late Q1 2026.",
        ],
    ),
    (
        "ESG and Sustainability",
        [
            "Scope 2 emissions down 17% on renewable PPAs.",
            "Diversity-of-board commitment maintained at 40% non-male directors.",
            "SASB reporting on track for inclusion in proxy.",
        ],
    ),
    (
        "Operational Risk Register",
        [
            "Vendor concentration risk flagged on two top-five suppliers.",
            "Business continuity plan tested in September; passed RTO targets.",
            "Insurance renewal completed at 6% premium increase.",
        ],
    ),
    (
        "Investor Relations Calendar",
        [
            "Q3 earnings call scheduled for November 4, 2025.",
            "Two sell-side conferences booked in December.",
            "Analyst day targeted for March 2026.",
        ],
    ),
    (
        "Closing — Board Asks",
        [
            "Approve FY2026 operating plan.",
            "Ratify cyber risk mitigation budget increase of $4.2M.",
            "Endorse Project Aurora term sheet upon receipt.",
        ],
    ),
]


def _build_deck(path: Path) -> None:
    """Materialise the deck described by SLIDES at ``path``.

    Uses python-pptx's blank-layout slide so the title placeholder
    behaves predictably across template versions.
    """
    # Imports gated on the module-level pytest.importorskip above.
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    # Layout 5 is "Title Only" on the default template — keeps the
    # title placeholder while letting us draw a freeform body box,
    # which avoids the auto-placeholder bullet styling weirdness.
    layout = prs.slide_layouts[5]

    for title_text, bullets in SLIDES:
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = title_text
        body_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.75), Inches(8.5), Inches(5.0))
        tf = body_box.text_frame
        tf.text = bullets[0]
        for bullet in bullets[1:]:
            para = tf.add_paragraph()
            para.text = bullet

    prs.save(str(path))


# ---------------------------------------------------------------------------
# Fixtures matching the K7 live test conventions
# ---------------------------------------------------------------------------


@pytest.fixture
def runtime() -> Any:
    # ``test_mode()`` installs an in-memory, globally-scoped VFS to
    # prevent session-memory leakage across pytest invocations (the
    # session_id below is hard-coded, which would let the default
    # disk VFS at ``.kaos-vfs`` false-green this test on re-run).
    from kaos_core.registry.container import KaosRuntime

    return KaosRuntime.test_mode()


@pytest.fixture
def context(runtime: Any) -> Any:
    from kaos_core.base.context import KaosContext

    return KaosContext.create(session_id="pptx-findings-live", runtime=runtime)


@pytest.fixture
def pptx_path() -> Any:
    """A freshly-built deck on disk, cleaned up after the test."""
    with tempfile.TemporaryDirectory() as tmp:
        path = Path(tmp) / "board_deck.pptx"
        _build_deck(path)
        yield path


async def _store_deck_artifact(runtime: Any, context: Any, path: Path) -> str:
    """Parse the PPTX + persist as a VFS artifact, returning its id."""
    from kaos_content.artifacts import store_document
    from kaos_office import parse_pptx

    doc = parse_pptx(str(path))
    manifest = await store_document(doc, runtime, context, name=path.stem)
    return manifest.artifact_id


# ---------------------------------------------------------------------------
# The test
# ---------------------------------------------------------------------------


@pytest.mark.live
@requires_anthropic
class TestPPTXFindingsLive:
    """K7 over a real PPTX board deck — composes kaos-office + K7."""

    async def test_token_selector_recovers_planted_cyber_mitigation(
        self,
        runtime: Any,
        context: Any,
        pptx_path: Path,
    ) -> None:
        """Full pipeline: PPTX → AST → VFS artifact → AgentFindingsTool.

        Token selector "cyber" should pin sentences on slide 8 of
        the deck. The synthesized answer must surface the planted
        mitigation phrase (proof the relevant sentence was retained
        through Phase 2 and cited in Phase 3, not hallucinated).
        """
        from kaos_agents.tools.findings import AgentFindingsTool

        artifact_id = await _store_deck_artifact(runtime, context, pptx_path)

        tool = AgentFindingsTool()
        result = await tool.execute(
            {
                "artifact_id": artifact_id,
                "question": (
                    "Which slide first mentions cyber risk, and what is the proposed mitigation?"
                ),
                "select_by": "token",
                "selector_arg": "cyber",
                "filter_model": critic_model(),
                "synthesis_model": critic_model(),
                "chunk_size": 20,
                "num_parallel": 3,
                "relevance_threshold": 0.4,
            },
            context,
        )

        assert not result.isError, f"Tool returned error: {result.text}"
        payload = result.structuredContent
        assert payload is not None, "Tool returned no structured payload"

        # Shape contract — same set the K7 wrapper test asserts.
        for key in (
            "artifact_id",
            "question",
            "answer",
            "findings",
            "total_enumerated",
            "total_filtered",
            "filter_calls",
            "filter_cost_usd",
            "synthesis_cost_usd",
            "total_cost_usd",
            "total_llm_calls",
        ):
            assert key in payload, f"missing {key} in payload"

        # Phase 1 — selector pinned at least one candidate. We plant
        # the word "cyber" on slides 8 and 15 (the closing-asks
        # slide also references the cyber budget), so floor of 1 is
        # safe and doesn't over-fit the wording.
        assert payload["total_enumerated"] >= 1, (
            f"Phase 1 enumerated {payload['total_enumerated']} candidates — "
            "selector regression or PPTX → sentence pipeline broken"
        )

        # Phase 2 — filter kept at least one survivor. If the
        # filter LLM drops everything that's a bug (the question
        # plainly asks about cyber risk; the planted sentence is
        # exactly that).
        assert payload["total_filtered"] >= 1, (
            "Phase 2 filter dropped every candidate; filter prompt "
            "or relevance_threshold too strict"
        )
        assert payload["total_filtered"] <= payload["total_enumerated"]

        # Phase 3 — synthesised answer mentions the planted
        # mitigation. We allow either substring to land — the LLM
        # may rephrase "multi-factor" or quote the verbatim
        # mitigation, both prove the right sentence reached
        # synthesis.
        answer_lc = payload["answer"].lower()
        assert "multi-factor" in answer_lc or "penetration testing" in answer_lc, (
            f"Synthesized answer doesn't surface the planted mitigation. "
            f"Answer: {payload['answer']!r}"
        )

        # Phase 3 — answer cites at least one finding_id. The
        # K6 synthesis prompt is explicitly grounded to cite the
        # short ids; the wrapper test confirms this on NDAs and we
        # mirror it here so a regression on PPTX provenance is
        # caught.
        from kaos_agents.patterns.findings import extract_finding_id_citations

        cited_ids = extract_finding_id_citations(payload["answer"])
        assert len(cited_ids) >= 1, (
            "Synthesis answer carries no finding_id citation; either "
            "the synthesis prompt regressed or no findings reached "
            "Phase 3"
        )

        # Cost envelope — single haiku per chunk + one haiku
        # synthesis on a 15-slide deck should comfortably fit under
        # $0.20. The K7 wrapper test caps at $0.50 because it uses
        # Sonnet for synthesis on a longer document; this test
        # uses Haiku everywhere so a tighter ceiling is appropriate.
        assert payload["filter_cost_usd"] > 0, "filter cost zero — usage wiring regression"
        assert payload["synthesis_cost_usd"] > 0, "synthesis cost zero — usage wiring regression"
        assert 0 < payload["total_cost_usd"] < 0.20, (
            f"total_cost_usd={payload['total_cost_usd']} outside sanity gate"
        )

        # AST provenance — every surviving finding has a block_ref
        # back into the kaos-content tree. This is what makes
        # downstream verify / source-quote tooling possible on
        # PPTX-derived findings.
        for finding in payload["findings"]:
            assert finding["finding_id"]
            assert isinstance(finding["text"], str)
            assert 0.0 <= finding["relevance"] <= 1.0
            assert finding["block_ref"] is not None, (
                "Finding missing block_ref — PPTX parser dropped provenance somewhere upstream"
            )
